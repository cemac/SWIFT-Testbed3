##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import glob
import matplotlib.image as mpimg
from datetime import datetime, timedelta
import argparse




# READ IN COMMAND LINE ARGUMENTS
dstring = ("Used to Generate PowerPoint from pngs")
hstring = ("--WG working group e.g. ENS \n --OUT file name string (typically contains meta data) \n --R region code 3 letter string")
parser = argparse.ArgumentParser(description=dstring)
parser.add_argument("--WG", help=hstring, type=str)
parser.add_argument("--OUT", help=hstring, type=str)
parser.add_argument("--R", help=hstring, type=str)
args = parser.parse_args()
if args.WG:
    WG = args.WG
else:
    WG = "na"

if args.OUT:
    OUTPUT_TAG = args.OUT
else:
    OUTPUT_TAG = "SWIFT_ppt"

if args.R:
    code = args.R
else:
    code = " "

# Define funtions
def country_decoder(code):
    countrynames = {"sen":"Senegal", "gha":"Ghana", "nga":"Nigeria",
                 "kya":"Kenya", "afr":"Africa",  "cafr":"Central Africa",
                 "eafr":"East Aftrica" ,"wafr":"West Africa"}
    return countrynames[code]

def city_decoder(code):
    citynames = {"DAK":"Dakar", "TBA":"Tambacounda", "TOU":"Touba",
                 "ACC":"Accra", "KUM":"Kumasi", "TAM":"Tamale", "ABU":"Abuja",
                 "KAN":"Kano", "LAG":"Lagos", "LAK":"Lake Victoria",
                 "MOM":"Mombasa Nairobi", "LAM":"Lamu", "VOI":"Voi",
                 "LOD":"Lodwar", "GAR":"Garissa", "MAN":"Mandera",
                 "MAR":"Marsabit", "KAK":"Kakamega", "KIT":"Kitale",
                 "KER":"Kericho", "KIS":"Kisumu", "NAI":"Nairobi", "NYE":"Nyeri",
                 "MER":"Meru", "NAK":"Nakuru", "NAR":"Narok", "MAC":"Machakos",
                 "KIT":'Kitui', "LAG":"Lagos", "POR":"Port_Harcourt","ENU":"Enugu",
                 "sen":"Senegal", "gha":"Ghana", "nga":"Nigeria", "kya":"Kenya",
                 "afr":"Africa",  "cafr":"Central Africa", "eafr":"East Aftrica",
                 "wafr":"West Africa"}
    return citynames[code]

def file_striper(filename):
    plottypes = ["nbhood_max", "stamp", "meteogram"]
    # Grab info from file
    filevars = filename.split("_")
    if "nbhood" in filevars:
        plot_type = "nbhood_max"
        tframe = filevars[1]
        date = filevars[9]
        hr = filevars[10]
        threshold = filevars[7]
        region = filevars[-4]
        L1 = int(filevars[-1].split(".")[0][1:])
        L2 = L1 + int(tframe[0:-2])
        date1 = datetime.strptime(str(date+hr[0:2]), '%Y%m%d%H')
        date2 = date1 + timedelta(hours=int(tframe[0:-2]))
        date1 = date1.strftime("%Y/%m/%d %H00Z")
        date2 = date2.strftime("%Y/%m/%d %H00Z")
        title = "Probability rainfall \n> " + threshold + " in " + tframe + "\n\n" + date1  + "-\n" + date2 + "\n(T+" + str(L1) + "-" + str(L2) + ")"
    elif "amount" in filevars:
        tframe = filevars[0]
        date = filevars[4]
        hr = filevars[5]
        L1 = int(filevars[-1].split(".")[0][1:])
        L2 = L1 + int(tframe[0:-2])
        region = filevars[3]
        date1 = datetime.strptime(str(date+hr[0:2]), '%Y%m%d%H')
        date2 = date1 + timedelta(hours=int(tframe[0:-2]))
        date1 = date1.strftime("%Y/%m/%d %H00Z")
        date2 = date2.strftime("%Y/%m/%d %H00Z")
        plot_type = "stamp"
        title = str(tframe)+ ' rainfall accumulation' + f"\n" + date1  + " - " + date2 + " (T+" + str(L1) + "-" + str(L2) + ")"
    elif "meteogram" in filevars:
        # Grab city code
        region = str(filevars[-1].split(".")[0])
        title = "Meteogram for\n" + str(city_decoder(region))
        plot_type = "meteogram"
    return title, plot_type, city_decoder(region)
# new
prs = pptx.Presentation()
prs.slide_height = 5143500
# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# set title
title = slide.shapes.title
title.text = str(country_decoder(code))
pic_left  = int(prs.slide_width * 0.15)
pic_top   = int(prs.slide_height * 0.01)

if WG == 'ENS':
    for g in glob.glob("*"):
        print(g)
        title, plot_type, region = file_striper(g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = mpimg.imread(g)
        # check aspect ratio and set width and height
        if img.shape[1] > img.shape[0]:  # w > h
            pic_width = int(prs.slide_width * 0.7)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
        else:
            pic_height = int(prs.slide_height * 0.98)
            pic_width = int(pic_height * img.shape[1] / img.shape[0])
        pic_left  = int((prs.slide_width - pic_width) * 0.5)
        pic_top  = int((prs.slide_height - pic_height) * 0.5)
        #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
        shapes = slide.shapes

        top  = Inches(0.2)
        if plot_type in ["nbhood_max", "meteogram"]:
            left = Inches(0.7)
            width = Inches(1)
            height = Inches(4)
            pic = slide.shapes.add_picture(g, int(pic_left*1.7), pic_top, pic_width, pic_height)
        else:
            left = Inches(2.5)
            width = Inches(4)
            height = Inches(0.5)
            pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

        txbox = slide.shapes.add_textbox(left, top, width, height)
        txbox.fill.solid()
        txbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = txbox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(16)


    prs.save("%s" % OUTPUT_TAG  + str(country_decoder(code)) + ".pptx")
else:
    for g in glob.glob("*"):
        print(g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shapes = slide.shapes
        shapes.title.text = g
        img = mpimg.imread(g)
        # check aspect ratio and set width and height
        if img.shape[1] > img.shape[0]:  # w > h
            pic_width = int(prs.slide_width * 0.7)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
        else:
            pic_height = int(prs.slide_height * 0.98)
            pic_width = int(pic_height * img.shape[1] / img.shape[0])
        pic_left  = int((prs.slide_width - pic_width) * 0.5)
        pic_top  = int((prs.slide_height - pic_height) * 0.5)
        #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
        pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
    prs.save("%s.pptx" % OUTPUT_TAG)
