##
# Auto PowerPoint genteration Tool
##
import pptx
import pptx.util
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import glob
import matplotlib.image as mpimg
from datetime import datetime, timedelta
import argparse


# READ IN COMMAND LINE ARGUMENTS
# There is an IF clause to moderate for WG Ensembles
dstring = ("Used to Generate PowerPoint from pngs")
hstring = ("--WG working group e.g. ENS \n --OUT file name string (typically contains meta data) \n --R region code 3 letter string")
parser = argparse.ArgumentParser(description=dstring)
parser.add_argument("--WG", help=hstring, type=str)
parser.add_argument("--OUT", help=hstring, type=str)
parser.add_argument("--R", help=hstring, type=str)
args = parser.parse_args()
# Set defaults or read in cmd line
if args.WG:
    WG = args.WG
else:
    WG = "na"

if args.OUT:
    OUTPUT_TAG = args.OUT
else:
    OUTPUT_TAG = "SWIFT_ppt"

if args.R:
    code = args.R
else:
    code = " "

# Define funtions


def country_decoder(code):
    countrynames = {"sen": "Senegal", "gha": "Ghana", "nga": "Nigeria",
                    "kya": "Kenya", "afr": "Africa",  "cafr": "Central Africa",
                    "eafr": "East_Africa", "wafr": "West_Africa"}
    return countrynames[code]


def city_decoder(code):
    citynames = {"DAK": "Dakar", "TBA": "Tambacounda", "TOU": "Touba",
                 "ACC": "Accra", "KUM": "Kumasi", "TAM": "Tamale", "ABU": "Abuja",
                 "KAN": "Kano", "LAG": "Lagos", "LAK": "Lake_Victoria",
                 "MOM": "Mombasa", "LAM": "Lamu", "VOI": "Voi",
                 "LOD": "Lodwar", "GAR": "Garissa", "MAN": "Mandera",
                 "MAR": "Marsabit", "KAK": "Kakamega", "KIT": "Kitale",
                 "KER": "Kericho", "KIM": "Kisumu", "NAI": "Nairobi", "NYE": "Nyeri",
                 "MER": "Meru", "NAK": "Nakuru", "NAR": "Narok", "MAC": "Machakos",
                 "KTU": 'Kitui', "POR": "Port_Harcourt", "ENU": "Enugu",
                 "sen": "Senegal", "gha": "Ghana", "nga": "Nigeria", "kya": "Kenya",
                 "afr": "Africa",  "cafr": "Central_Africa", "eafr": "East_Africa",
                 "wafr": "West_Africa","GHANA_ABE": "Abetifi ", "GHANA_ACCACA": "Accaca academy",
                 "GHANA_ACC": "Accaca Girls", "GHANA_ADA": "Ada ",
                 "GHANA_AGISS": "GHANA_AGISS", "GHANA_AKAN": "Akanseringa",
                 "GHANA_AKA": "Akatsi", "GHANA_AKU": " ",
                 "GHANA_AODA": "Akim Oda ", "GHANA_ARCH": " ",
                 "GHANA_ATTC": "Attc circle", "GHANA_AXI": "Axim",
                 "GHANA_BAA": "Baasa", "GHANA_BOL": "Bole",
                 "GHANA_GUL": "Gulumpe", "GHANA_HO": "Ho",
                 "GHANA_JIR": "Jirapa", "GHANA_KKRA": "Kete Krachi",
                 "GHANA_KOF": "Koforidua", "GHANA_KOLEBU": "GHANA_KOLEBU",
                 "GHANA_KSI": "GHANA KSI", "GHANA_KUS": "Kushegu",
                 "GHANA_LOA": "Loagri NO2 ", "GHANA_MAMO": "Maobi Poly Clinic",
                 "GHANA_MPE": "Mempeasem", "GHANA_MUL": "Mulpido",
                 "GHANA_NAK": "Nakpaboni", "GHANA_NAV": "Navrongo",
                 "GHANA_PRESEC": "Presec-OSU", "GHANA_SAL": "Saltpond",
                 "GHANA_SAW": "Sawla", "GHANA_SBEK": "Sefwi-Bekwai",
                 "GHANA_SNY": "Sunyani", "GHANA_SUM": "Sumbrungu",
                 "GHANA_TDI": "Takoradi", "GHANA_TEM": "Temale",
                 "GHANA_TLE": "GHANA_TLE", "GHANA_TOL": "Tolon",
                 "GHANA_VAR": "Varenpare", "GHANA_WA": "WA",
                 "GHANA_WEIJA": "Weija-Dam ", "GHANA_WEN": "Wenchi",
                 "GHANA_WESLEY": "Wesley Academy", "GHANA_YEN": "Yendi",}
    return citynames[code]


def file_striper(filename):
    #
    # Grab info from file and split by _
    plottypes = ["nbhood_max", "stamp", "meteogram"]
    # Grab info from file
    filevars = filename.split("_")
    if "nbhood" in filevars:
        # set plot type
        plot_type = "nbhood_max"
        # grab metadata
        tframe = filevars[1]
        date = filevars[9]
        hr = filevars[10]
        threshold = filevars[7]
        region = filevars[-4]
        # Extract lead time
        L1 = int(filevars[-1].split(".")[0][1:])
        L2 = L1 - int(tframe[0:-2])
        # Turn into date time and get date range and return to date string
        date1 = datetime.strptime(str(date + hr[0:2]), '%Y%m%d%H')
        date2 = date1 - timedelta(hours=int(tframe[0:-2]))
        date1 = date1.strftime("%Y/%m/%d %H00Z")
        date2 = date2.strftime("%Y/%m/%d %H00Z")
        # Generate title
        title = "Probability rainfall \n> " + threshold + " in " + tframe + \
            "\n\n" + date2 + "-\n" + date1 + \
                "\n(T+" + str(L2) + "-" + str(L1) + ")"
    elif "amount" in filevars:
        # pick out stamp plots
        # grab metadata
        tframe = filevars[0]
        date = filevars[4]
        hr = filevars[5]
        # Extract lead time
        L1 = int(filevars[-1].split(".")[0][1:])
        L2 = L1 - int(tframe[0:-2])
        region = filevars[3]
        # Turn into date time and get date range and return to date string
        date1 = datetime.strptime(str(date + hr[0:2]), '%Y%m%d%H')
        date2 = date1 - timedelta(hours=int(tframe[0:-2]))
        date1 = date1.strftime("%Y/%m/%d %H00Z")
        date2 = date2.strftime("%Y/%m/%d %H00Z")
        plot_type = "stamp"
        # Generate title
        title = str(tframe) + ' rainfall accumulation' + f"\n" + \
            date2 + " - " + date1 + " (T+" + str(L2) + "-" + str(L1) + ")"
    elif "meteogram" in filevars:
        # Grab city code
        region = str(filevars[-1].split(".")[0])
        if str(filevars[-2]) == "GHANA":
            region = str(filevars[-2])+'_'+region
        # Generate title
        title = "Meteogram for\n" + str(city_decoder(region))
        plot_type = "meteogram"
    return title, plot_type, city_decoder(region)


# new PowerPoint
prs = pptx.Presentation()
prs.slide_height = 5143500
# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# set title
title = slide.shapes.title
title.text = str(country_decoder(code))
pic_left = int(prs.slide_width * 0.15)
pic_top = int(prs.slide_height * 0.01)

if WG == 'ENS':
    # For Ensembles
    # Loop through pngs
    for g in glob.glob("*"):
        try:
            title, plot_type, region = file_striper(g)
        except:
            continue
        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Read image
        img = mpimg.imread(g)
        # check aspect ratio and set width and height
        if img.shape[1] > img.shape[0]:  # w > h
            pic_width = int(prs.slide_width * 0.7)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
        else:
            pic_height = int(prs.slide_height * 0.98)
            pic_width = int(pic_height * img.shape[1] / img.shape[0])
        pic_left = int((prs.slide_width - pic_width) * 0.5)
        pic_top = int((prs.slide_height - pic_height) * 0.5)
        # Get shapes
        shapes = slide.shapes
        # For nbhood and meteogram the title text is to the left

        if plot_type in ["nbhood_max", "meteogram"]:
            top = Inches(0.2)
            left = Inches(0.7)
            width = Inches(1)
            height = Inches(4)
            # shift pic a long a bit to make room for text
            pic = slide.shapes.add_picture(g, int(pic_left * 1.7),
                                           pic_top, pic_width, pic_height)
        else:
            # for the other plots cover title of plot with text
            # The Region plots are a different size to the country plots
            if region in ["Africa", "East_Africa", "Central_Africa", "West_Africa"]:
                top = Inches(0.5)
                if region == "Africa":
                    # The Africa plot is a different size to
                    # the other region plots
                    top = Inches(1)
            else:
                top = Inches(0.2)
            left = Inches(2.5)
            width = Inches(4)
            height = Inches(0.5)
            pic = slide.shapes.add_picture(
                g, pic_left, pic_top, pic_width, pic_height)
        # create a text box
        txbox = slide.shapes.add_textbox(left, top, width, height)
        # white fill it
        txbox.fill.solid()
        txbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = txbox.text_frame
        p = tf.add_paragraph()
        # add title from metadata and set fontsize
        p.text = title
        p.font.size = Pt(16)

    # save with country name
    prs.save("%s" % OUTPUT_TAG + str(country_decoder(code)) + ".pptx")
else:
    # For synoptic/other plots
    for g in glob.glob("*"):
        print(g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = mpimg.imread(g)
        # check aspect ratio and set width and height
        if img.shape[1] > img.shape[0]:  # w > h
            pic_width = int(prs.slide_width * 0.7)
            pic_height = int(pic_width * img.shape[0] / img.shape[1])
        else:
            pic_height = int(prs.slide_height * 0.98)
            pic_width = int(pic_height * img.shape[1] / img.shape[0])
        pic_left  = int((prs.slide_width - pic_width) * 0.5)
        pic_top  = int((prs.slide_height - pic_height) * 0.5)
        #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
        pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
        if WG == 'SYNOP':
            # Add legend based on filename
            # e.g. 20210831_0600_069_WA_convective.png
            filevars = g.split(".")[0].split("_")
            region = filevars[3]
            chart_type = filevars[4]
            if chart_type == "synthesis":
                chart_type = "summary"
            elif chart_type == "wa-jets-waves":
                chart_type = "jets"
            path = f'../../legends/{region}_{chart_type}.png'
            try:
                img = mpimg.imread(path)
                # set width and position to fit beside chart
                w = pic_left
                h = int(w * img.shape[0] / img.shape[1])
                if h > prs.slide_height:
                    h = int(prs.slide_height * 0.98)
                    w = int(h * img.shape[1] / img.shape[0])
                left  = 0
                top  = int((prs.slide_height - h) * 0.5)
                legend = slide.shapes.add_picture(path, left, top, width=w, height=h)
            except:
                print("Couldn't find legend file: ", path)
    if WG == "SYNOP":
        if OUTPUT_TAG.find("nowcasting") > -1:
            title.text += " simplified synthetic analysis charts for the nowcasting briefings"
    prs.save("%s.pptx" % OUTPUT_TAG)
